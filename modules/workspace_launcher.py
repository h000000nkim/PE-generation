"""
워크스페이스 생성 + Claude Code 터미널 실행
- Notion에서 파싱한 데이터를 파일로 저장 (코드 자동화 영역)
- Claude Code가 수행할 분석/작성 지침을 CLAUDE.md로 작성
- macOS Terminal.app에서 Claude Code 실행
"""

import os
import re
import json
import logging
import subprocess
import requests
from pathlib import Path

logger = logging.getLogger(__name__)

BASE_DIR = Path(__file__).resolve().parent.parent / "workspaces"
LOCKED_FILE = BASE_DIR / ".locked_tasks.json"

# block_id → workspace path 매핑 (런타임 캐시)
_workspace_map: dict[str, Path] = {}
_workspace_prefix_map: dict[str, Path] = {}  # block_id 앞8자 → path

# 서버 시작 시 워크스페이스 디렉토리 사전 매핑
if BASE_DIR.exists():
    for _d in BASE_DIR.iterdir():
        if _d.is_dir() and not _d.name.startswith("."):
            _parts = _d.name.rsplit("_", 1)
            if len(_parts) == 2 and len(_parts[1]) >= 8:
                _workspace_prefix_map[_parts[1][:8]] = _d


def get_locked_ids() -> set:
    """잠금된 과제 ID 목록"""
    if LOCKED_FILE.exists():
        return set(json.loads(LOCKED_FILE.read_text(encoding="utf-8")))
    return set()


def set_locked_ids(ids: set):
    """잠금 과제 ID 저장"""
    LOCKED_FILE.parent.mkdir(parents=True, exist_ok=True)
    LOCKED_FILE.write_text(json.dumps(list(ids), ensure_ascii=False), encoding="utf-8")


def _safe_dirname(name: str, block_id: str) -> str:
    safe = re.sub(r'[^\w가-힣ㄱ-ㅎㅏ-ㅣ _-]', '_', name or "unknown")
    return f"{safe}_{block_id[:8]}"


def _download_file(url: str, dest: Path):
    try:
        r = requests.get(url, timeout=30, stream=True)
        r.raise_for_status()
        with open(dest, "wb") as f:
            for chunk in r.iter_content(8192):
                f.write(chunk)
        return True
    except Exception:
        return False


def create_workspace(task: dict) -> Path:
    """
    코드 자동화 영역: STEP 0 + STEP 1 + STEP 2(다운로드) + 학생 데이터 수집
    - 디렉토리 생성
    - 첨부파일/가이드/생기부 다운로드
    - 파싱된 데이터를 구조화된 파일로 저장
    - Claude Code용 CLAUDE.md 작성
    """
    block_id = task.get("block_id", "unknown")
    dirname = _safe_dirname(task.get("name", ""), block_id)
    ws = BASE_DIR / dirname
    ws.mkdir(parents=True, exist_ok=True)

    # 경로 매핑 즉시 등록
    _workspace_map[block_id] = ws
    _workspace_prefix_map[block_id[:8]] = ws

    # block_id → workspace 매핑 저장
    _workspace_map[block_id] = ws

    # ── MCP 설정 복사 (Word MCP 서버 등) ──
    import shutil
    project_mcp = Path(__file__).resolve().parent.parent / ".mcp.json"
    if project_mcp.exists():
        shutil.copy2(project_mcp, ws / ".mcp.json")

    # ── Claude Code 설정 복사 (MCP 자동 승인 등) ──
    claude_dir = ws / ".claude"
    claude_dir.mkdir(exist_ok=True)
    settings = {
        "permissions": {
            "allow": [
                "Bash(*)",
                "Read(*)",
                "Write(*)",
                "Edit(*)",
                "Glob(*)",
                "Grep(*)",
                "WebFetch(*)",
                "WebSearch(*)",
                "mcp__word-document-server__*",
                "mcp__hwpx-document-server__*",
                "mcp__claude_ai_Canva__*"
            ]
        },
        "enabledMcpjsonServers": [
            "word-document-server",
            "hwpx-document-server"
        ],
        "enableAllProjectMcpServers": True
    }
    (claude_dir / "settings.local.json").write_text(
        json.dumps(settings, indent=2, ensure_ascii=False), encoding="utf-8"
    )

    # ── 파일 다운로드 (코드 자동화) ──
    files_dir = ws / "files"
    files_dir.mkdir(exist_ok=True)

    downloaded_attachments = []
    downloaded_guides = []
    downloaded_bio = []

    for att in task.get("attachments", []):
        if att.get("url"):
            fname = att.get("name", "attachment")
            if _download_file(att["url"], files_dir / fname):
                downloaded_attachments.append(fname)

    for gf in task.get("guide_files", []):
        if gf.get("url"):
            fname = gf.get("name", "guide")
            if _download_file(gf["url"], files_dir / fname):
                downloaded_guides.append(fname)

    for bf in task.get("bio_files", []):
        if bf.get("url"):
            fname = bf.get("name", "bio")
            if _download_file(bf["url"], files_dir / fname):
                downloaded_bio.append(fname)

    # ── 문서 텍스트 사전 추출 (.pptx, .docx, .pdf → .txt) ──
    all_downloaded = downloaded_attachments + downloaded_guides + downloaded_bio
    for fname in all_downloaded:
        fpath = files_dir / fname
        txt_path = files_dir / f"{fname}.txt"
        try:
            if fname.lower().endswith(".pptx"):
                from pptx import Presentation
                prs = Presentation(str(fpath))
                lines = []
                for i, slide in enumerate(prs.slides, 1):
                    lines.append(f"=== 슬라이드 {i} ===")
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = para.text.strip()
                                if t:
                                    lines.append(t)
                    lines.append("")
                txt_path.write_text("\n".join(lines), encoding="utf-8")
            elif fname.lower().endswith(".docx"):
                from docx import Document as DocxDocument
                doc = DocxDocument(str(fpath))
                lines = [p.text for p in doc.paragraphs if p.text.strip()]
                txt_path.write_text("\n".join(lines), encoding="utf-8")
            elif fname.lower().endswith(".pdf"):
                import fitz
                doc = fitz.open(str(fpath))
                text_lines = []
                has_text = False
                for i, page in enumerate(doc, 1):
                    page_text = page.get_text().strip()
                    text_lines.append(f"=== 페이지 {i} ===")
                    text_lines.append(page_text)
                    text_lines.append("")
                    if len(page_text) > 20:
                        has_text = True
                # 텍스트가 있으면 .txt로 저장
                if has_text:
                    txt_path.write_text("\n".join(text_lines), encoding="utf-8")
                else:
                    # 스캔 PDF → 페이지별 이미지로 변환
                    for i, page in enumerate(doc, 1):
                        pix = page.get_pixmap(dpi=150)
                        img_path = files_dir / f"{fname}_p{i}.png"
                        pix.save(str(img_path))
                    # 안내 파일 생성
                    txt_path.write_text(
                        f"스캔 PDF — 텍스트 없음. 페이지별 이미지로 변환됨:\n"
                        + "\n".join(f"  - {fname}_p{i}.png" for i in range(1, len(doc) + 1)),
                        encoding="utf-8"
                    )
                doc.close()
            elif fname.lower().endswith(".hwpx"):
                import zipfile
                with zipfile.ZipFile(str(fpath)) as zf:
                    texts = []
                    # section XML에서 텍스트 추출
                    for zname in zf.namelist():
                        if re.match(r'Contents/section\d+\.xml$', zname):
                            xml = zf.read(zname).decode("utf-8")
                            for m in re.finditer(r'<hp:t[^>]*>([^<]+)</hp:t>', xml):
                                texts.append(m.group(1))
                    # PrvText.txt 대체
                    if not texts and "Preview/PrvText.txt" in zf.namelist():
                        t = zf.read("Preview/PrvText.txt").decode("utf-8", errors="ignore").strip()
                        if t:
                            texts.append(t)
                    if texts:
                        txt_path.write_text("\n".join(texts), encoding="utf-8")
                    else:
                        # 미리보기 이미지라도 추출
                        if "Preview/PrvImage.png" in zf.namelist():
                            img_path = files_dir / f"{fname}_preview.png"
                            img_path.write_bytes(zf.read("Preview/PrvImage.png"))
                            txt_path.write_text(
                                f"HWPX — 텍스트 없음. 미리보기 이미지로 변환됨:\n  - {fname}_preview.png",
                                encoding="utf-8"
                            )
        except Exception as e:
            logger.warning(f"[workspace] 텍스트 추출 실패 ({fname}): {e}")

    # ── 과거 과제 이력 정리 (코드 자동화) ──
    past_tasks_text = ""
    if task.get("past_tasks"):
        lines = []
        for pt in task["past_tasks"]:
            line = f"- **{pt.get('title', '(제목없음)')}** | 과목: {pt.get('subject', '—')} | 상태: {pt.get('status', '—')}"
            if pt.get('semester'):
                line += f" | {pt['semester']}"
            if pt.get('due_date'):
                line += f" | 마감: {pt['due_date']}"
            if pt.get('request'):
                line += f"\n  요청 핵심: {pt['request'][:200]}"
            lines.append(line)
        past_tasks_text = "\n".join(lines)

    # ── 파일 목록 텍스트 ──
    att_list = "\n".join(f"  - files/{f} (평가기준/양식)" for f in downloaded_attachments) if downloaded_attachments else "  (없음)"
    guide_list = "\n".join(f"  - files/{f} (완성물/참고용 — 분석 불필요)" for f in downloaded_guides) if downloaded_guides else "  (없음)"
    bio_list = "\n".join(f"  - files/{f} (생기부)" for f in downloaded_bio) if downloaded_bio else "  (없음)"

    # ── CLAUDE.md 작성 (학생/과제 정보 + 작업 절차만 — 전역 규칙은 상위 CLAUDE.md에서 자동 로드) ──
    claude_md = f"""# 워크스페이스 — {task.get('title', '—')}

## 역할
너는 수행평가 작문 및 기획 전문가다. 아래 데이터와 files/ 폴더의 첨부파일을 바탕으로 작업을 수행한다.

---

## [파싱 완료] 학생 정보
- 이름: {task.get('name', '—')}
- 학교: {task.get('school', '—')}
- 학년: {task.get('grade', '—')}
- 진로/계열: {task.get('major', '—')}
{f"- 목표 학과: {task['target_dept']}" if task.get('target_dept') else ""}
{f"- 소속: {task['affiliation']}" if task.get('affiliation') else ""}
- 학생 코드: {task.get('student_code', '—')}
{f"- 연락처: {task.get('phone', '')}" if task.get('phone') else ""}

## [파싱 완료] 과제 정보
- 과목(mGaa): {task.get('subject', '—')}
- 학기(wuL:): {task.get('semester', '—')}
- 수행평가 형식(Dogm): {task.get('submit_type', '—')}
- 컨설팅 종류(CrVV): {task.get('activity', '—')}
- 진행 상태(owtr): {task.get('status', '—')}
{f"- 신청일: {task['apply_date']}" if task.get('apply_date') else ""}
{f"- 학교 제출일: {task['due_date']}" if task.get('due_date') else ""}
{f"- 키워드(UfkI): {task['keyword']}" if task.get('keyword') else ""}
{f"- 비고: {task['note']}" if task.get('note') else ""}

## [파싱 완료] 요청사항 (RSt|)
{task.get('request_msg') or '없음'}

{f"## [파싱 완료] 세부능력 및 특기사항 (RuwY)" + chr(10) + task['setech'] if task.get('setech') else ""}

{f"## [파싱 완료] 생기부 방향성 (NVS@)" + chr(10) + task['bio_direction'] if task.get('bio_direction') else ""}

{f"## [파싱 완료] 시험범위 / 수업내용 (HR~K)" + chr(10) + task['study_range'] if task.get('study_range') else ""}

## [파싱 완료] 과거 수행 과제 이력
{past_tasks_text if past_tasks_text else "(이력 없음)"}

## [다운로드 완료] 첨부 파일
평가기준 및 양식 (gAMl):
{att_list}
수행평가 가이드 (mPfP) — 완성물/참고용, 분석 대상 아님:
{guide_list}
생기부 파일:
{bio_list}

---

## 작업 절차 (Claude Code가 수행할 작업)

STEP 1, STEP 2 다운로드는 이미 완료됨. 아래부터 수행한다.

### STEP 2. 이미지/파일 분석
files/ 폴더의 모든 첨부파일을 열어서 분석한다.
- **모든 문서 파일(.pptx, .docx, .pdf)의 텍스트는 이미 추출 완료되어 같은 폴더에 `.txt` 파일로 존재한다**
  - 예: `파일명.pptx` → `파일명.pptx.txt`, `생기부.pdf` → `생기부.pdf.txt`
  - 반드시 `.txt` 파일을 Read로 읽어서 분석할 것
  - 원본 .pptx/.docx/.pdf를 직접 파싱하려고 시도하지 말 것 (패키지 설치 금지: python-pptx, python-docx, pymupdf, poppler 등)
  - `.txt` 파일이 이미 있으므로 별도 도구나 패키지가 필요 없다
- 이미지 파일(.jpeg, .png 등)은 직접 열어서 분석한다
- 평가기준, 양식 항목, 배점, 분량, 지시사항을 빠짐없이 파악한다
- 이미지 첨부파일이 여러 장인 경우 **전부** 분석할 것 (일부만 보고 판단하지 말 것)
- 구글폼 양식이 있는 경우 폼 항목을 빠짐없이 파악할 것

### STEP 3. 사람 분석 (먼저 출력)
위 파싱된 데이터를 바탕으로 아래 항목을 정리해서 출력한다:
- 이름, 학교, 학년(위 grade 값 그대로), 진로/계열
- 생기부 방향성
- 탐구 주제 이력 요약
- 이전 과업 목록 (과목 + 상태 + 요청 핵심)
- 학생의 특징, 관심사, 작업 스타일

### STEP 4. 과제 분석 (이미지 + 텍스트 종합)
아래 항목을 모두 확인하고 출력한다:
- 과목, 수행평가 형식, 제출 마감
- 평가기준 및 양식(gAMl) 이미지에서 파악한 배점/항목
- (수행평가 가이드 mPfP는 완성물이므로 분석 불필요)
- 요청사항 분석 — 선생님 지시, 압박 질문 예상
- 키워드 참고
- 이전 활동과의 연결점
- **제출 형식 확인**: PPT 포함 여부, 발표 여부, 구글폼 여부, Word 여부 등

### STEP 5. 과제 작성
STEP 3, 4의 분석 결과를 출력한 후, 바로 작성에 들어간다.
작성 원칙, PPT/Canva 지침, Word 서식 규칙, 자주 발생하는 실수, 응답 규칙은 상위 CLAUDE.md(rules/)를 따른다.
"""

    (ws / "CLAUDE.md").write_text(claude_md, encoding="utf-8")

    return ws


def launch_pre_analysis(task: dict, workspace: Path) -> bool:
    """
    Claude Code 터미널로 STEP 2~4 사전분석 실행.
    분석 결과는 analysis.md + analysis.json에 저장됨.
    """
    instruction = f"""CLAUDE.md를 읽고, files/ 폴더의 첨부파일을 분석한 뒤,
**STEP 2(파일 분석) → STEP 3(사람 분석) → STEP 4(과제 분석)**만 수행하세요.
STEP 5(과제 작성)는 하지 마세요.

과제: {task.get('title', '')}
과목: {task.get('subject', '')}
형식: {task.get('submit_type', '')}
학생: {task.get('name', '')} ({task.get('grade', '')}, {task.get('school', '')})

분석 결과를 아래 형식으로 출력하고, analysis.md 파일로 저장하세요:

[STEP 3 — 사람 분석]
(이름, 학교, 학년, 진로/계열, 생기부 방향성, 탐구 주제 이력, 이전 과업 이력, 특징 및 작업 스타일)

[STEP 4 — 과제 분석]
(과목, 형식, 마감, 평가기준/배점, 요청사항 분석, 제출 형식 확인, 모호한 점 지적)

[작업 계획]
(어떤 순서로, 어떤 내용을, 어떤 형식으로 작성할지 구체적 계획)

마지막으로 analysis.json 파일도 저장하세요:
{{"status": "ok", "analysis": "(analysis.md 내용 전체)"}}"""

    block_id = task.get("block_id", "")
    _save_memo_log(block_id, "사전분석", "STEP 2~4 사전분석 실행")
    return launch_background(workspace, instruction, block_id, "사전분석")


def get_analysis(block_id: str) -> dict | None:
    """워크스페이스에서 캐싱된 사전분석 읽기"""
    ws = get_workspace_path(block_id)
    if not ws:
        return None
    analysis_file = ws / "analysis.json"
    if analysis_file.exists():
        return json.loads(analysis_file.read_text(encoding="utf-8"))
    return None


def delete_analysis(block_id: str) -> bool:
    """사전분석 결과 삭제 (analysis.json + analysis.md)"""
    ws = get_workspace_path(block_id)
    if not ws:
        return False
    deleted = False
    for fname in ("analysis.json", "analysis.md"):
        f = ws / fname
        if f.exists():
            f.unlink()
            deleted = True
    return deleted


def _save_memo_log(block_id: str, action: str, memo: str):
    """추가 지시사항/수정 요청 이력 저장"""
    ws = get_workspace_path(block_id)
    if not ws:
        return
    log_file = ws / "memo_log.json"
    logs = []
    if log_file.exists():
        logs = json.loads(log_file.read_text(encoding="utf-8"))
    logs.append({
        "action": action,
        "memo": memo,
        "timestamp": _time.strftime("%Y-%m-%d %H:%M:%S"),
        "result_snapshot": None,  # job 완료 시 채워짐
    })
    log_file.write_text(json.dumps(logs, ensure_ascii=False, indent=2), encoding="utf-8")


def _attach_result_to_log(block_id: str, warning: str | None = None):
    """마지막 로그 항목에 현재 result.json 스냅샷 + 오류 정보 + 변경 파일 첨부"""
    ws = get_workspace_path(block_id)
    if not ws:
        return
    log_file = ws / "memo_log.json"
    if not log_file.exists():
        return
    try:
        logs = json.loads(log_file.read_text(encoding="utf-8"))
        if not logs:
            return

        entry = logs[-1]

        # 오류 기록
        if warning:
            entry["error"] = warning
            entry["result_snapshot"] = entry.get("result_snapshot") or []
        else:
            entry["error"] = None

            if entry.get("action") == "사전분석":
                # 사전분석은 analysis.json을 스냅샷으로 첨부
                analysis_file = ws / "analysis.json"
                if analysis_file.exists():
                    entry["result_snapshot"] = [{"label": "사전분석 완료", "type": "text", "file": "analysis.md"}]
            else:
                # 정상 완료 시 result.json 스냅샷 첨부
                result_file = ws / "result.json"
                if result_file.exists():
                    result = json.loads(result_file.read_text(encoding="utf-8"))
                    summary = []
                    for o in result.get("outputs", []):
                        s = {"label": o.get("label", ""), "type": o.get("type", ""), "file": o.get("file", "")}
                        if o.get("canva_edit_url"):
                            s["canva_edit_url"] = o["canva_edit_url"]
                        if o.get("pptx_download_url"):
                            s["pptx_download_url"] = o["pptx_download_url"]
                        summary.append(s)
                    entry["result_snapshot"] = summary

                # 수정 요청인 경우: 이 수정으로 변경된 파일 목록 기록
                if entry.get("action") == "수정 요청":
                    changed = _detect_changed_files(ws, entry.get("timestamp", ""))
                    if changed:
                        entry["changed_files"] = changed

        log_file.write_text(json.dumps(logs, ensure_ascii=False, indent=2), encoding="utf-8")
    except Exception:
        pass


def _detect_changed_files(ws: Path, since_timestamp: str) -> list[str]:
    """수정 요청 타임스탬프 이후에 변경된 산출물 파일 목록"""
    import datetime
    changed = []
    try:
        since = datetime.datetime.strptime(since_timestamp, "%Y-%m-%d %H:%M:%S")
    except (ValueError, TypeError):
        return []

    for f in ws.iterdir():
        if not f.is_file():
            continue
        if f.suffix in (".md", ".docx", ".hwpx", ".json") and f.name not in (
            "CLAUDE.md", "analysis.md", "analysis.json", "memo_log.json",
            "verification.json", ".warning.json"
        ) and not f.name.startswith("."):
            try:
                mtime = datetime.datetime.fromtimestamp(f.stat().st_mtime)
                if mtime > since:
                    changed.append(f.name)
            except Exception:
                pass
    return changed


def attach_verification_to_log(block_id: str):
    """최신 검증 결과 전체를 가장 최근 비-사전분석 로그 항목에 첨부"""
    ws = get_workspace_path(block_id)
    if not ws:
        return
    log_file = ws / "memo_log.json"
    vf = ws / "verification.json"
    if not log_file.exists() or not vf.exists():
        return
    try:
        logs = json.loads(log_file.read_text(encoding="utf-8"))
        if not logs:
            return
        verification = json.loads(vf.read_text(encoding="utf-8"))
        # 사전분석이 아닌 가장 최근 항목을 찾아서 첨부
        target = None
        for entry in reversed(logs):
            if entry.get("action") != "사전분석":
                target = entry
                break
        if target is None:
            target = logs[-1]  # 사전분석밖에 없으면 마지막에라도 첨부
        target["verification"] = verification
        log_file.write_text(json.dumps(logs, ensure_ascii=False, indent=2), encoding="utf-8")
        logger.info(f"[verify-log] {block_id} 검증 결과 → '{target.get('action')}' ({target.get('timestamp')}) 항목에 첨부")
    except Exception as e:
        logger.error(f"[verify-log] {block_id} 첨부 실패: {e}")


def get_memo_log(block_id: str) -> list:
    """추가 지시사항 이력 조회"""
    ws = get_workspace_path(block_id)
    if not ws:
        return []
    log_file = ws / "memo_log.json"
    if log_file.exists():
        return json.loads(log_file.read_text(encoding="utf-8"))
    return []


def get_verification_status(block_id: str) -> str | None:
    """검증 상태 조회 — pass/fail/warning/None"""
    ws = get_workspace_path(block_id)
    if not ws:
        return None
    vf = ws / "verification.json"
    if not vf.exists():
        return None
    try:
        data = json.loads(vf.read_text(encoding="utf-8"))
        return data.get("status")
    except Exception:
        return None


def run_verification(task: dict) -> bool:
    """
    산출물 검증 실행 — Claude Code 백그라운드로 검증 수행.
    result.json/result.md를 읽고, CLAUDE.md의 기준 대비 검증.
    결과는 verification.json에 저장됨.
    """
    block_id = task.get("block_id", "")
    ws = get_workspace_path(block_id)
    if not ws:
        return False

    result = get_result(block_id)
    if not result:
        return False

    instruction = f"""CLAUDE.md를 읽고, 현재 산출물(result.json, result.md, 및 관련 파일)을 검증하세요.

## 검증 항목

### 1. 기본 요건 충족 확인
- 제출 형식(Dogm)과 실제 산출물 형식이 일치하는지
- 요청사항(RSt|)에 명시된 사항이 모두 반영되었는지
- 평가기준 이미지의 채점 항목이 모두 충족되었는지

### 2. 내용 품질 검증
- AI 특유의 문체(나열식, "~할 수 있다" 반복 등)가 남아있는지
- 학생의 진로/계열과 내용이 자연스럽게 연결되는지
- 분량이 적절한지 (평가기준 기준)
- 소제목/번호 매기기 남용이 없는지 (양식 요구가 아닌 경우)

### 3. 사실 관계 검증
- 출처 URL이 있는 경우 실제 접속 가능한지 확인
- 학년, 학교 등 기본 정보가 정확한지

### 4. 파일 무결성
- result.json의 outputs 배열과 실제 파일 존재 여부 일치 확인
- docx 파일이 정상적으로 열리는지 (파일 크기 > 0)
- PPT가 있는 경우 canva_edit_url, pptx_download_url 존재 확인

## 출력 형식

검증 결과를 verification.json으로 저장:
```json
{{
  "status": "pass" 또는 "fail" 또는 "warning",
  "timestamp": "검증 시각",
  "score": 0-100,
  "summary": "전체 요약 (1-2줄)",
  "checks": [
    {{
      "category": "기본 요건 | 내용 품질 | 사실 관계 | 파일 무결성",
      "item": "검증 항목명",
      "status": "pass | fail | warning",
      "detail": "상세 설명"
    }}
  ],
  "recommendations": ["개선 제안 1", "개선 제안 2"]
}}
```

과제: {task.get('title', '')}
과목: {task.get('subject', '')}
형식: {task.get('submit_type', '')}
학생: {task.get('name', '')} ({task.get('grade', '')}, {task.get('school', '')})"""

    return launch_background(ws, instruction, block_id, "검증")


def build_instruction(task: dict, user_memo: str = "") -> str:
    """Claude Code에 전달할 초기 프롬프트 — 캐싱된 분석이 있으면 STEP 5부터 시작"""
    block_id = task.get("block_id", "")
    _save_memo_log(block_id, "초안 작성", user_memo or "(추가 지시 없음)")
    _refresh_claude_md(task)
    analysis = get_analysis(block_id)

    # 파일명 생성용 정보
    subj = (task.get('subject', '') or '').replace(' ', '')
    name = task.get('name', '')
    filename_hint = f"예시: {subj}_주제_{name}.docx 또는 {subj}_주제_{name}.hwpx" if subj and name else ""

    if analysis and analysis.get("status") == "ok":
        # 사전분석 완료 → STEP 5만 실행
        instruction = f"""CLAUDE.md를 읽고, files/ 폴더의 첨부파일을 확인한 뒤,
아래 사전분석 결과를 바탕으로 **STEP 5(과제 작성)**를 바로 수행해주세요.

⚠️⚠️ 최우선 규칙 — 파일명:
- 산출물 파일명은 반드시 `과목_과제제목_이름.확장자` 형식으로 저장할 것
- {filename_hint}
- **result.docx, result.md, result.hwpx 등 임시 이름 절대 금지**
- 사전분석의 "저장 경로" 항목에 result.docx라 적혀 있어도 무시하고 위 규칙을 따를 것

## 사전분석 결과 (STEP 3 + STEP 4 완료)
{analysis['analysis']}

## 과제 기본정보
과제: {task.get('title', '')}
과목: {task.get('subject', '')}
형식: {task.get('submit_type', '')}
학생: {task.get('name', '')} ({task.get('grade', '')}, {task.get('school', '')})

{f"## 추가 지시사항" + chr(10) + user_memo if user_memo else ""}

산출물별로 개별 파일 + result.json(outputs 배열)으로 저장해주세요."""
    else:
        # 사전분석 없음 → 전체 수행
        instruction = f"""CLAUDE.md를 읽고, files/ 폴더의 첨부파일을 모두 분석한 뒤,
STEP 2 → STEP 3 → STEP 4 → STEP 5 순서로 작업을 수행해주세요.

⚠️⚠️ 최우선 규칙 — 파일명:
- 산출물 파일명은 반드시 `과목_과제제목_이름.확장자` 형식으로 저장할 것
- {filename_hint}
- **result.docx, result.md, result.hwpx 등 임시 이름 절대 금지**

과제: {task.get('title', '')}
과목: {task.get('subject', '')}
형식: {task.get('submit_type', '')}
학생: {task.get('name', '')} ({task.get('grade', '')}, {task.get('school', '')})

{f"추가 지시사항: {user_memo}" if user_memo else ""}

산출물별로 개별 파일 + result.json(outputs 배열)으로 저장해주세요."""
    return instruction


def _refresh_claude_md(task: dict):
    """기존 워크스페이스의 CLAUDE.md만 최신 템플릿으로 재생성 (파일 다운로드 안 함)"""
    ws = get_workspace_path(task.get("block_id", ""))
    if not ws:
        return
    # create_workspace를 호출하면 같은 경로에 CLAUDE.md를 덮어씀
    # 이미 다운로드된 파일은 _download_file이 기존 파일 위에 다시 쓰지만 빠름
    create_workspace(task)


def build_revision_instruction(task: dict, revision_memo: str) -> str:
    """기존 산출물 수정 프롬프트"""
    block_id = task.get("block_id", "")
    _refresh_claude_md(task)
    ws = get_workspace_path(block_id)

    # 현재 산출물 목록 파악
    existing_files = []
    if ws:
        for f in ws.iterdir():
            if f.is_file() and f.suffix in (".md", ".docx", ".hwpx", ".txt") and f.name not in ("CLAUDE.md", "analysis.md", ".prompt.txt"):
                existing_files.append(f.name)
        result_file = ws / "result.json"
        if result_file.exists():
            result_data = json.loads(result_file.read_text(encoding="utf-8"))
            outputs = result_data.get("outputs", [])
            for o in outputs:
                if o.get("file"):
                    existing_files.append(o["file"])

    file_list = "\n".join(f"  - {f}" for f in sorted(set(existing_files))) if existing_files else "  (산출물 없음)"

    return f"""CLAUDE.md를 읽고, 기존 산출물을 수정해주세요.

## 현재 산출물 파일
{file_list}

## 수정 요청
{revision_memo}

## 수정 원칙
- 기존 파일을 읽고 요청된 부분만 수정할 것
- 수정하지 않은 부분은 원문 그대로 유지할 것
- 수정 후 같은 파일명으로 덮어쓸 것
- result.json도 업데이트할 것
- 사용자가 직접 확정한 문장/표현은 수정하지 말 것

과제: {task.get('title', '')}
학생: {task.get('name', '')} ({task.get('grade', '')}, {task.get('school', '')})"""


def launch_revision(task: dict, revision_memo: str) -> bool:
    """기존 산출물 수정을 위한 Claude Code 백그라운드 실행"""
    block_id = task.get("block_id", "")
    if revision_memo:
        _save_memo_log(block_id, "수정 요청", revision_memo)
    ws = get_workspace_path(block_id)
    if not ws:
        return False
    instruction = build_revision_instruction(task, revision_memo)
    return launch_background(ws, instruction, block_id, "수정")


# 실행 중인 프로세스 추적: block_id → [job, job, ...] (다중 작업 지원)
import time as _time
import threading as _threading_mod

_running_jobs: dict[str, list] = {}
_jobs_lock = _threading_mod.Lock()


def _kill_same_label_jobs(block_id: str, label: str):
    """같은 과제의 같은 작업 유형(label) 중 실행 중인 프로세스를 종료한다.
    다른 유형의 작업(예: 수정과 분석)은 병렬로 유지된다."""
    with _jobs_lock:
        jobs = _running_jobs.get(block_id, [])
        for job in jobs:
            if job["label"] == label and job["status"] == "running":
                proc = job.get("process")
                if proc:
                    try:
                        proc.terminate()
                        proc.wait(timeout=5)
                    except Exception:
                        try:
                            proc.kill()
                        except Exception:
                            pass
                    job["status"] = "cancelled"
                    job["end_time"] = _time.time()
                    logger.info(f"[background] 기존 {label} 작업 종료 (pid={job['pid']}, block={block_id})")


def launch_background(workspace_path: Path, instruction: str, block_id: str, label: str = "작업") -> bool:
    """Claude Code를 백그라운드 프로세스로 실행 (-p 모드).
    같은 과제 + 같은 작업 유형이 이미 실행 중이면 기존 세션을 종료하고 새로 시작한다."""
    # 같은 label의 기존 작업 종료 (다른 label은 유지 → 병렬 허용)
    _kill_same_label_jobs(block_id, label)

    ws = str(workspace_path)

    # 작업별 고유 파일 (덮어쓰기 방지)
    job_id = f"{label}_{int(_time.time())}"
    prompt_file = workspace_path / f".prompt_{job_id}.txt"
    prompt_file.write_text(instruction, encoding="utf-8")

    output_file = workspace_path / f".output_{job_id}.txt"
    error_file = workspace_path / f".error_{job_id}.txt"

    stdin_f = stdout_f = stderr_f = None
    try:
        stdin_f = open(prompt_file, "r")
        stdout_f = open(output_file, "w")
        stderr_f = open(error_file, "w")

        proc = subprocess.Popen(
            [
                "claude", "-p",
                "--output-format", "text",
                "--permission-mode", "bypassPermissions",
            ],
            stdin=stdin_f,
            stdout=stdout_f,
            stderr=stderr_f,
            cwd=ws,
        )

        job = {
            "job_id": job_id,
            "pid": proc.pid,
            "process": proc,
            "start_time": _time.time(),
            "status": "running",
            "label": label,
            "workspace": ws,
            "output_file": str(output_file),
            "error_file": str(error_file),
            "warning": None,
        }

        with _jobs_lock:
            if block_id not in _running_jobs:
                _running_jobs[block_id] = []
            _running_jobs[block_id].append(job)

        # 백그라운드 스레드에서 완료 감시 + 문제 감지 + 임시파일 정리
        import threading
        def _watch():
            try:
                proc.wait()
            finally:
                # Popen 후 파일 핸들을 스레드에서 안전하게 닫기
                for fh in (stdin_f, stdout_f, stderr_f):
                    try:
                        fh.close()
                    except Exception:
                        pass
            job["status"] = "complete" if proc.returncode == 0 else "error"
            job["end_time"] = _time.time()
            job["returncode"] = proc.returncode
            # output 파일에서 문제 감지
            warning = _detect_issues(output_file, error_file)
            job["warning"] = warning
            if warning:
                _save_warning(ws, job_id, label, warning)
            else:
                _clear_warning(ws)
            # 검증이 아닌 경우만 결과물 이력 기록
            if label != "검증":
                _attach_result_to_log(block_id, warning)
            # 완료된 job의 임시 파일 정리 (최신 1개만 유지)
            _cleanup_job_files(Path(ws))
            # 초안 작성 또는 수정 성공 시 자동 검증 트리거
            if label in ("초안 작성", "수정") and proc.returncode == 0 and not warning:
                _auto_verify(block_id)
            # 검증 완료 시 결과를 지시이력에 첨부
            if label == "검증" and proc.returncode == 0:
                attach_verification_to_log(block_id)
                logger.info(f"[verify] {block_id} 검증 결과 이력 첨부 완료")

        threading.Thread(target=_watch, daemon=True).start()
        return True

    except Exception as e:
        # 예외 시 파일 핸들 정리 + 프로세스 종료
        for fh in (stdin_f, stdout_f, stderr_f):
            if fh:
                try:
                    fh.close()
                except Exception:
                    pass
        logger.error(f"[background] 실행 실패: {e}")
        return False


def _save_warning(workspace: str, job_id: str, label: str, warning: str):
    """경고를 워크스페이스 파일에 저장"""
    wpath = Path(workspace) / ".warning.json"
    wpath.write_text(json.dumps({
        "job_id": job_id,
        "label": label,
        "warning": warning,
        "timestamp": _time.strftime("%Y-%m-%d %H:%M:%S"),
    }, ensure_ascii=False), encoding="utf-8")


def _clear_warning(workspace: str):
    """경고 파일 제거 (정상 완료 시)"""
    wpath = Path(workspace) / ".warning.json"
    if wpath.exists():
        wpath.unlink()


def get_warning(block_id: str) -> dict | None:
    """워크스페이스의 미해결 경고 조회"""
    ws = get_workspace_path(block_id)
    if not ws:
        return None
    wpath = ws / ".warning.json"
    if wpath.exists():
        return json.loads(wpath.read_text(encoding="utf-8"))
    return None


def _auto_verify(block_id: str):
    """초안 작성 완료 후 자동 검증 트리거"""
    try:
        # result.json이 있는지 확인
        ws = get_workspace_path(block_id)
        if not ws:
            return
        result_file = ws / "result.json"
        if not result_file.exists():
            return
        # notion_parser를 여기서 import (순환 import 방지)
        from modules.notion_parser import parse_task_from_block
        task = parse_task_from_block(block_id)
        run_verification(task)
        logger.info(f"[auto-verify] {task.get('title', block_id)} 자동 검증 시작")
    except Exception as e:
        logger.error(f"[auto-verify] {block_id} 실패: {e}")


def _cleanup_job_files(workspace: Path):
    """완료된 job 임시 파일 정리 — 최신 1개만 유지"""
    for prefix in [".output_", ".error_", ".prompt_", ".expect_"]:
        files = sorted(workspace.glob(f"{prefix}*"), key=lambda f: f.stat().st_mtime, reverse=True)
        # 최신 1개 유지, 나머지 삭제
        for f in files[1:]:
            try:
                f.unlink()
            except Exception:
                pass


def _detect_issues(output_file: Path, error_file: Path) -> str | None:
    """output/error 파일에서 권한 문제 등 감지"""
    issues = []
    for fpath in [output_file, error_file]:
        fpath = Path(fpath)
        if not fpath.exists():
            continue
        try:
            text = fpath.read_text(encoding="utf-8", errors="ignore")[:3000]
        except Exception:
            continue
        # 권한/승인 문제
        if any(kw in text for kw in ["권한", "permission", "승인", "Permission denied", "not allowed"]):
            if "Canva" in text or "canva" in text:
                issues.append("Canva MCP 권한 문제")
            elif "word" in text.lower() or "Word" in text:
                issues.append("Word MCP 권한 문제")
            elif "hwpx" in text.lower() or "HWPX" in text:
                issues.append("HWPX MCP 권한 문제")
            else:
                issues.append("MCP 권한 문제")
        # 인증 만료
        if any(kw in text for kw in ["auth", "token expired", "401", "인증"]):
            issues.append("인증 만료 가능성")
        # rate limit
        if any(kw in text for kw in ["rate limit", "429", "too many requests"]):
            issues.append("API 요청 제한")
    return " · ".join(issues) if issues else None


def get_job_status(block_id: str) -> dict | None:
    """실행 중인 작업 상태 조회 — 모든 활성 작업 반환"""
    with _jobs_lock:
        jobs = _running_jobs.get(block_id, [])
        if not jobs:
            return None

        active = []
        warnings = []
        for j in jobs:
            elapsed = _time.time() - j["start_time"]
            active.append({
                "job_id": j["job_id"],
                "status": j["status"],
                "label": j["label"],
                "elapsed_seconds": int(elapsed),
                "pid": j["pid"],
                "warning": j.get("warning"),
            })
            if j.get("warning"):
                warnings.append(j["warning"])

        # 완료된 지 5분 지난 작업 정리
        _running_jobs[block_id] = [
            j for j in jobs
            if j["status"] == "running" or _time.time() - j.get("end_time", _time.time()) < 300
        ]

    warning_text = " · ".join(set(warnings)) if warnings else None

    # running이 있으면 running 우선 반환
    running = [a for a in active if a["status"] == "running"]
    if running:
        return {"status": "running", "jobs": active, "label": ", ".join(r["label"] for r in running),
                "elapsed_seconds": max(a["elapsed_seconds"] for a in running), "pid": running[0]["pid"],
                "warning": warning_text}

    # 전부 완료면 가장 최근 것
    latest = max(active, key=lambda a: a["elapsed_seconds"])
    return {"status": latest["status"], "jobs": active, "label": latest["label"],
            "elapsed_seconds": latest["elapsed_seconds"], "pid": latest["pid"],
            "warning": warning_text}


def get_workspace_path(block_id: str) -> Path | None:
    """block_id로 워크스페이스 경로 조회 (매핑 캐시 → 사전 매핑 → 디스크 탐색)"""
    if block_id in _workspace_map:
        return _workspace_map[block_id]
    # 사전 매핑에서 즉시 찾기
    short = block_id[:8]
    if short in _workspace_prefix_map:
        path = _workspace_prefix_map[short]
        _workspace_map[block_id] = path
        return path
    # 폴백: 디스크 탐색 (새로 생성된 워크스페이스)
    if BASE_DIR.exists():
        for d in BASE_DIR.iterdir():
            if d.is_dir() and short in d.name:
                _workspace_map[block_id] = d
                _workspace_prefix_map[short] = d
                return d
    return None


def get_result(block_id: str) -> dict | None:
    """워크스페이스에서 result.json 읽기"""
    ws = get_workspace_path(block_id)
    if not ws:
        return None
    result_file = ws / "result.json"
    if result_file.exists():
        return json.loads(result_file.read_text(encoding="utf-8"))
    # result.json 없으면 result.md라도 읽기
    md_file = ws / "result.md"
    if md_file.exists():
        return {"text": md_file.read_text(encoding="utf-8")}
    return None


def save_result(block_id: str, data: dict) -> bool:
    """워크스페이스에 result.json 저장"""
    ws = get_workspace_path(block_id)
    if not ws:
        return False
    result_file = ws / "result.json"
    # 기존 결과와 병합
    existing = {}
    if result_file.exists():
        existing = json.loads(result_file.read_text(encoding="utf-8"))
    existing.update(data)
    result_file.write_text(json.dumps(existing, ensure_ascii=False, indent=2), encoding="utf-8")
    return True
